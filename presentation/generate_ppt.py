# -*- coding: utf-8 -*-
"""
발표자료(PPTX) 생성기 — 에타 강의평 RAG 프로젝트 (고퀄리티 테마 + 데모 스크린샷).
실행: python generate_ppt.py  →  강의평RAG_발표자료.pptx
"""
import os
import struct

from pptx import Presentation
from pptx.util import Pt, Emu, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))
SHOT = os.path.join(HERE, "screenshots")
OUT = os.path.join(HERE, "강의평RAG_발표자료_고퀄.pptx")

FONT = "맑은 고딕"
NAVY = RGBColor(0x16, 0x21, 0x3E)
BLUE = RGBColor(0x2D, 0x5B, 0xFF)
RED = RGBColor(0xE0, 0x32, 0x2E)
GRAY = RGBColor(0x4A, 0x4A, 0x4A)
LIGHT = RGBColor(0xF2, 0xF4, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

EMU = 914400
SW, SH = 13.333, 7.5


def png_size(path):
    with open(path, "rb") as f:
        head = f.read(24)
    if head[:8] == b"\x89PNG\r\n\x1a\n":
        w, h = struct.unpack(">II", head[16:24])
        return w, h
    return 1600, 1000


def _set(run, size, color, bold=False):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def rect(slide, x, y, w, h, color, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    return tf


def fit_image(slide, path, x, y, w, h):
    iw, ih = png_size(path)
    scale = min(w / (iw / 96.0), h / (ih / 96.0))   # px→inch(96dpi 가정) 후 박스에 맞춤
    dw, dh = (iw / 96.0) * scale, (ih / 96.0) * scale
    px, py = x + (w - dw) / 2, y + (h - dh) / 2
    # 흰 카드 + 테두리
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(px - 0.08), Inches(py - 0.08),
                                  Inches(dw + 0.16), Inches(dh + 0.16))
    card.fill.solid(); card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(0xD0, 0xD6, 0xE0); card.line.width = Pt(1)
    card.shadow.inherit = False
    slide.shapes.add_picture(path, Inches(px), Inches(py), Inches(dw), Inches(dh))


def base(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, SW, SH, WHITE)
    return s


def header(s, title, page):
    rect(s, 0, 0, SW, 1.15, NAVY)
    rect(s, 0, 1.15, SW, 0.06, RED)
    rect(s, 0.55, 0.34, 0.16, 0.5, RED)
    tf = textbox(s, 0.85, 0.28, 11.5, 0.7, MSO_ANCHOR.MIDDLE)
    r = tf.paragraphs[0].add_run(); r.text = title; _set(r, 27, WHITE, True)
    # footer
    rect(s, 0, SH - 0.04, SW, 0.04, LIGHT)
    ft = textbox(s, 0.5, SH - 0.45, 12.3, 0.35, MSO_ANCHOR.MIDDLE)
    p = ft.paragraphs[0]
    r = p.add_run(); r.text = "강의평 RAG AI 비서"; _set(r, 9, GRAY)
    p2 = ft.paragraphs[0]; p2.alignment = PP_ALIGN.LEFT
    pg = textbox(s, 11.8, SH - 0.45, 1.0, 0.35, MSO_ANCHOR.MIDDLE)
    pr = pg.paragraphs[0]; pr.alignment = PP_ALIGN.RIGHT
    rr = pr.add_run(); rr.text = str(page); _set(rr, 10, GRAY)


def bullets(tf, items, size=17):
    for i, b in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        sub = b.startswith("  ")
        p.space_after = Pt(9)
        run = p.add_run()
        run.text = ("◦  " + b.strip()) if sub else ("▪  " + b)
        _set(run, size - 3 if sub else size, GRAY if sub else NAVY, bold=False)


def title_slide(prs):
    s = base(prs)
    rect(s, 0, 0, SW, SH, NAVY)
    rect(s, 0, 4.55, SW, 0.07, RED)
    tf = textbox(s, 1.0, 1.9, 11.3, 2.4, MSO_ANCHOR.MIDDLE)
    for i, line in enumerate(["에브리타임 강의평 기반",
                              "RAG 수강신청 AI 비서"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = line; _set(r, 40, WHITE, True)
    sub = textbox(s, 1.0, 4.8, 11.3, 1.6, MSO_ANCHOR.TOP)
    for i, line in enumerate([
            "우리 학교 실제 강의평을 근거로, 자연어 한 줄에 맞춤 수강 추천",
            "가천대 5과목 · 직접 크롤링 187건 · 벡터 RAG · Claude"]):
        p = sub.paragraphs[0] if i == 0 else sub.add_paragraph()
        p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(6)
        r = p.add_run(); r.text = line
        _set(r, 15, RGBColor(0xC8, 0xD2, 0xE8))


def content_slide(prs, page, title, items, image=None, caption=None):
    s = base(prs)
    header(s, title, page)
    if image and os.path.exists(image):
        tf = textbox(s, 0.85, 1.55, 6.1, 5.4)
        bullets(tf, items, size=16)
        fit_image(s, image, 7.1, 1.5, 5.8, 5.0)
        if caption:
            cap = textbox(s, 7.1, 6.55, 5.8, 0.4, MSO_ANCHOR.MIDDLE)
            cp = cap.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
            r = cp.add_run(); r.text = caption; _set(r, 11, GRAY)
    else:
        tf = textbox(s, 0.95, 1.6, 11.4, 5.3)
        bullets(tf, items, size=18)


def image_slide(prs, page, title, image, caption):
    s = base(prs)
    header(s, title, page)
    fit_image(s, image, 1.6, 1.5, 10.1, 5.2)
    cap = textbox(s, 1.0, 6.75, 11.3, 0.45, MSO_ANCHOR.MIDDLE)
    cp = cap.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
    r = cp.add_run(); r.text = caption; _set(r, 12, GRAY)


def end_slide(prs):
    s = base(prs)
    rect(s, 0, 0, SW, SH, NAVY)
    tf = textbox(s, 1.0, 2.7, 11.3, 1.4, MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "감사합니다"; _set(r, 40, WHITE, True)
    sub = textbox(s, 1.0, 4.2, 11.3, 0.8, MSO_ANCHOR.MIDDLE)
    p = sub.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "github.com/whtjdgjs12/everytime"
    _set(r, 15, RGBColor(0xC8, 0xD2, 0xE8))


def main():
    prs = Presentation()
    prs.slide_width = Inches(SW); prs.slide_height = Inches(SH)
    home = os.path.join(SHOT, "01_home.png")
    answer = os.path.join(SHOT, "02_answer.png")
    chat = os.path.join(SHOT, "03_chat.png")

    title_slide(prs)
    content_slide(prs, 2, "문제 (Pain Points)", [
        "수강신청마다 강의평 수십~수백 개를 직접 스크롤·비교 → 시간 낭비",
        "'학점 3.5 사수', '과제 적은 과목' 같은 목표 맞춤 추천이 없음",
        "일반 AI 챗봇은 우리 학교 강의평을 몰라 거짓 정보(환각)를 답함",
        "→ 우리 학교 데이터에 근거해 답하는 지능형 비서가 필요",
    ])
    content_slide(prs, 3, "해결: RAG (검색 증강 생성)", [
        "RAG = '오픈북 시험' 방식 AI",
        "① 검색: 강의평 DB에서 질문 관련 리뷰를 찾아옴",
        "② 생성: 찾아온 리뷰'만' 근거로 답을 생성",
        "모델 재학습이 아님 → 데이터만 바꾸면 즉시 우리 학교용",
        "근거 없으면 '데이터 없음'이라 솔직히 거부 → 환각 방지",
    ])
    content_slide(prs, 4, "서비스 개요", [
        "자연어 한 줄 질문 → 근거(출처) 포함 답변",
        "과목 필터 · 참고 리뷰 수 조절 · 출처 펼쳐보기",
        "API 키 없이도 동작(오프라인 요약)",
        "키 넣으면 Claude로 자연어 답변 업그레이드",
    ], image=home, caption="▲ 실제 채팅 웹앱 화면")
    content_slide(prs, 5, "데이터: 직접 수집한 실제 강의평", [
        "대상: 가천대 5과목",
        "  확률과 통계 · 모바일프로그래밍 · 빅데이터분석개론",
        "  스마트기기시스템 · 블록체인개론",
        "수집: Playwright 브라우저 자동화 → 강의평 187건",
        "스키마: 학교·교수·과목·수강학기·평점·리뷰",
        "정제: 중복 제거, 외국인반/공학인증/영어강의 제외, 과목명 통일",
    ])
    content_slide(prs, 6, "시스템 아키텍처", [
        "[크롤러] everytime_crawler.py (Playwright)",
        "  → reviews_real.csv (187건)",
        "[인덱싱] 임베딩(LSA) → 벡터 저장소(numpy)",
        "[검색] 질문 임베딩 → 코사인 최근접 top-k + 교수/과목 필터",
        "[생성] Claude API (없으면 규칙기반 오프라인 폴백)",
        "[서비스] Streamlit 채팅 웹앱 (공개 URL)",
    ])
    content_slide(prs, 7, "핵심 기능 (요구사항 3.1~3.4)", [
        "3.1 의미 검색: 글자가 안 겹쳐도 뜻으로 관련 강의평 검색",
        "3.2 속성 필터: 교수·과목으로 좁혀 정확도 향상",
        "3.3 출처 표기: (가천대 모바일프로그래밍 26년 1학기 평점 5점 리뷰 참고)",
        "3.4 환각 방지: 데이터 없는 교수/과목은 솔직히 거부",
    ])
    image_slide(prs, 8, "데모 — 질문과 근거 기반 답변", answer,
                "▲ '모바일프로그래밍 학점 잘 주는 교수?' → 교수별 비교 + 출처(학기·평점) 자동 표기")
    content_slide(prs, 9, "기술 의사결정 · 트러블슈팅", [
        "ChromaDB·onnx·torch가 환경에서 네이티브 충돌(세그폴트)",
        "  → 순수 numpy 벡터 저장소로 동일 기능 구현(의존성 0)",
        "에타는 Vue SPA: 요소 렌더 대기 + 강의평 탭 URL 직접 이동",
        "검색→클릭 방식의 멈춤 → 강의 URL 수집 후 직접 이동으로 해결",
        "약관·계정 보호: 자동수집 리스크 인지, 최소 요청·지연 적용",
    ])
    content_slide(prs, 10, "답변 생성 · 비용 · 배포", [
        "생성: Claude(기본 Haiku 4.5) / 키 없으면 오프라인 요약",
        "프롬프트 캐싱으로 반복 입력 토큰 비용 절감",
        "비용(대략): 질문 1건 ≈ 6원(Haiku) — 하루 수십 질문도 월 수천 원",
        "배포: Streamlit Cloud 공개 URL, API 키는 Secrets로 안전 관리",
    ])
    content_slide(prs, 11, "기대 효과 & 향후 계획", [
        "수강신청 탐색 시간 대폭 단축, 목표(학점/과제) 맞춤 수강",
        "근거(출처) 기반이라 신뢰도 높고 환각 없음",
        "확장: 과목 추가 크롤링, 강의계획서(PDF)·시험정보 보강",
        "캠퍼스 RAG 생태계(족보·장학·학사공지)로 확장 가능",
    ])
    end_slide(prs)

    prs.save(OUT)
    print("저장:", OUT, "| 슬라이드", len(prs.slides._sldIdLst))


if __name__ == "__main__":
    main()
